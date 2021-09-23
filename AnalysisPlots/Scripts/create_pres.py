from pptx import Presentation 
from pptx.util import Inches 
import sys, os 

# Image Adding Functions 
def add_centered_image(prs, lyt, name, img_path):
    slide = prs.slides.add_slide(lyt)
    title = slide.shapes.title 
    title.text = name
    center_left = Inches(2)
    center_top = Inches(2) 
    img = slide.shapes.add_picture(img_path, center_left, center_top)

def add_two_image(prs, lyt, name, img_path):
    slide = prs.slides.add_slide(lyt)
    title = slide.shapes.title 
    title.text = name
    left_1 = Inches(0)
    left_2 = Inches(5) 
    top = Inches(2) 

    img = slide.shapes.add_picture(img_path[0], left_1, top, Inches(5), Inches(5))
    img = slide.shapes.add_picture(img_path[1], left_2, top, Inches(5), Inches(5)) 


# Defining necessary paths/vars 
try: 
    dataset = sys.argv[1].strip()
    smoothing = sys.argv[2].strip()
except Exception as e:
    print("Please input dataset and smoothing as arguments\n", e) 
    exit()

save_path = f'../Data/{dataset}/'
fig_path = save_path + 'Figures/' 

# Creating presentation 
prs = Presentation() 

# Creating title slide 
lyt=prs.slide_layouts[0] # choosing a slide layout
slide=prs.slides.add_slide(lyt) # adding a slide
title=slide.shapes.title # assigning a title
subtitle=slide.placeholders[1] # placeholder for subtitle
title.text= f'{dataset}' # title
subtitle.text="Made with Python :) By Samir Chowdhury" # subtitle

# Adding images  
lyt = prs.slide_layouts[5]
title_beg = f'{dataset}-{smoothing} '
file_name_beg = f'{dataset}_smoothing_{smoothing}_'
center_left = Inches(2)
center_top = Inches(2) 

# Compare Behavior Analysis
cb_path = fig_path + 'compare_behavior_analysis/'

figs = ['phase_portrait', 
        'transform_validation_hilbert',
        'detection_histogram_hilbert', 
        'waveforms_hilbert',
        'pca_hilbert']

for fig in figs: 
    path = cb_path + file_name_beg + fig + '.png'
    title = title_beg + fig
    add_centered_image(prs, lyt, title, path)

# Compare Pressure Epochs
pe_path = fig_path + 'comp_pressure_epoch/'

# cycle_dur_plot 
path = pe_path + file_name_beg + 'cycle_dur_plot.png'
title = title_beg + 'cycle duration' 
add_centered_image(prs, lyt, title, path) 

# Epoch plots 
epochs_path = pe_path + 'Epochs/'

# Counting number of epochs 
num_epochs = 0

for base, dirs, files in os.walk(epochs_path):
    for directories in dirs:
        num_epochs += 1 

# Iterating through each folder and adding plots 
for epoch in range(1, num_epochs+1):
    name = f'epoch_{epoch}'
    epoch_folder_path = epochs_path + name + '/'
    epoch_data_path = epoch_folder_path + file_name_beg + name + '_data.png'
    epoch_pca_path = epoch_folder_path + file_name_beg + name + '_pca.png'
    add_two_image(prs, lyt, title_beg + name, [epoch_data_path, epoch_pca_path])

# Saving the presentation 
prs.save(save_path + f"{dataset}.pptx") 

